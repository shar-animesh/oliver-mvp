"""Content-addressed attachment storage and bounded document extraction."""

from __future__ import annotations

import base64
import binascii
import hashlib
import os
from io import BytesIO
from pathlib import Path
from typing import Protocol
from uuid import uuid4
from zipfile import ZipFile

from azure.core.exceptions import ResourceExistsError
from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobServiceClient
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy import select
from sqlalchemy.orm import Session

from utils.models.api import EmailAttachmentInput
from utils.postgres import AttachmentBlobDb, EmailAttachmentDb, EmailMessageDb

_SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}
_ARCHIVE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
_MAX_ARCHIVE_ENTRIES = 10_000
_MAX_PDF_PAGES = 500
_MAX_DOCUMENT_BLOCKS = 50_000
_MAX_SPREADSHEET_ROWS = 100_000
_MAX_PRESENTATION_SLIDES = 2_000


class AttachmentValidationError(ValueError):
    """An attachment cannot be accepted safely by the configured pipeline."""


class AttachmentStore(Protocol):
    """Persistence port for immutable attachment bytes."""

    def put(self, content_hash: str, content: bytes) -> str:
        """Persist bytes once and return their storage URI."""


class FileSystemAttachmentStore:
    """Local content-addressed store used during development and tests."""

    def __init__(self, root: str) -> None:
        self._root = Path(root).resolve()

    def put(self, content_hash: str, content: bytes) -> str:
        target_directory = self._root / content_hash[:2]
        target_directory.mkdir(parents=True, exist_ok=True)
        target = target_directory / content_hash
        if not target.exists():
            temporary = target.with_name(f"{content_hash}.{uuid4().hex}.tmp")
            try:
                temporary.write_bytes(content)
                os.replace(temporary, target)
            finally:
                temporary.unlink(missing_ok=True)
        return target.as_uri()


class AzureBlobAttachmentStore:
    """Production content-addressed store authenticated with workload identity."""

    def __init__(self, *, account_url: str, container: str) -> None:
        credential = DefaultAzureCredential(exclude_interactive_browser_credential=True)
        self._container = BlobServiceClient(account_url=account_url, credential=credential).get_container_client(container)

    def put(self, content_hash: str, content: bytes) -> str:
        blob = self._container.get_blob_client(f"sha256/{content_hash[:2]}/{content_hash}")
        try:
            blob.upload_blob(content, overwrite=False)
        except ResourceExistsError:
            pass
        return blob.url


def _bounded(value: str, maximum: int) -> str:
    return value[:maximum]


def _extract_pdf(content: bytes) -> tuple[str, dict[str, object]]:
    reader = PdfReader(BytesIO(content))
    if len(reader.pages) > _MAX_PDF_PAGES:
        raise AttachmentValidationError(f"PDF contains more than {_MAX_PDF_PAGES} pages")
    return "\n\n".join(page.extract_text() or "" for page in reader.pages), {"pages": len(reader.pages)}


def _extract_docx(content: bytes) -> tuple[str, dict[str, object]]:
    document = Document(BytesIO(content))
    block_count = len(document.paragraphs) + sum(len(table.rows) for table in document.tables)
    if block_count > _MAX_DOCUMENT_BLOCKS:
        raise AttachmentValidationError(f"Word document contains more than {_MAX_DOCUMENT_BLOCKS} text blocks")
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    table_rows = [" | ".join(cell.text for cell in row.cells) for table in document.tables for row in table.rows]
    return "\n".join([*paragraphs, *table_rows]), {"paragraphs": len(paragraphs), "tables": len(document.tables)}


def _extract_xlsx(content: bytes) -> tuple[str, dict[str, object]]:
    workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
    try:
        lines: list[str] = []
        row_count = 0
        worksheet_count = len(workbook.worksheets)
        for worksheet in workbook.worksheets:
            lines.append(f"Sheet: {worksheet.title}")
            for row in worksheet.iter_rows(values_only=True):
                row_count += 1
                if row_count > _MAX_SPREADSHEET_ROWS:
                    raise AttachmentValidationError(f"Excel workbook contains more than {_MAX_SPREADSHEET_ROWS} rows")
                values = [str(value) for value in row if value is not None]
                if values:
                    lines.append(" | ".join(values))
        return "\n".join(lines), {"worksheets": worksheet_count}
    finally:
        workbook.close()


def _extract_pptx(content: bytes) -> tuple[str, dict[str, object]]:
    presentation = Presentation(BytesIO(content))
    if len(presentation.slides) > _MAX_PRESENTATION_SLIDES:
        raise AttachmentValidationError(f"PowerPoint contains more than {_MAX_PRESENTATION_SLIDES} slides")
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if slide_text:
            lines.append(f"Slide {index}\n" + "\n".join(slide_text))
    return "\n\n".join(lines), {"slides": len(presentation.slides)}


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
}


def _validate_archive_size(content: bytes, *, maximum: int) -> None:
    """Reject Office archives that would expand beyond extraction limits."""
    with ZipFile(BytesIO(content)) as archive:
        entries = archive.infolist()
        if len(entries) > _MAX_ARCHIVE_ENTRIES:
            raise AttachmentValidationError(f"Office document contains more than {_MAX_ARCHIVE_ENTRIES} archive entries")
        total_uncompressed = sum(entry.file_size for entry in entries)
        if total_uncompressed > maximum:
            raise AttachmentValidationError(f"Office document expands to {total_uncompressed} bytes; maximum is {maximum}")


class AttachmentService:
    """Validate, store, extract and link attachments without duplicating payloads."""

    def __init__(self, *, store: AttachmentStore, max_bytes: int, max_uncompressed_bytes: int, max_extracted_chars: int) -> None:
        self._store = store
        self._max_bytes = max_bytes
        self._max_uncompressed_bytes = max_uncompressed_bytes
        self._max_extracted_chars = max_extracted_chars

    def ingest(self, database: Session, message: EmailMessageDb, attachment: EmailAttachmentInput) -> EmailAttachmentDb:
        existing = database.scalar(
            select(EmailAttachmentDb).where(
                EmailAttachmentDb.message_id == message.id,
                EmailAttachmentDb.provider_attachment_id == attachment.attachment_id,
            )
        )
        if existing is not None:
            return existing

        try:
            content = base64.b64decode(attachment.content_base64, validate=True)
        except (binascii.Error, ValueError) as error:
            raise AttachmentValidationError(f"Attachment {attachment.file_name!r} is not valid base64") from error
        if len(content) > self._max_bytes:
            raise AttachmentValidationError(f"Attachment {attachment.file_name!r} is {len(content)} bytes; maximum is {self._max_bytes}")

        file_name = Path(attachment.file_name).name
        extension = Path(file_name).suffix.lower()
        content_hash = hashlib.sha256(content).hexdigest()
        blob = database.get(AttachmentBlobDb, content_hash)
        if blob is None:
            storage_uri = self._store.put(content_hash, content)
            blob = AttachmentBlobDb(
                content_hash=content_hash,
                storage_uri=storage_uri,
                size_bytes=len(content),
                content_type=attachment.content_type,
            )
            extractor = _EXTRACTORS.get(extension)
            if extension not in _SUPPORTED_EXTENSIONS or extractor is None:
                blob.extraction_status = "UNSUPPORTED"
                blob.extraction_error = f"Unsupported attachment type {extension or '(none)'}"
            else:
                try:
                    if extension in _ARCHIVE_EXTENSIONS:
                        _validate_archive_size(content, maximum=self._max_uncompressed_bytes)
                    extracted_text, metadata = extractor(content)
                    blob.extracted_text = _bounded(extracted_text.strip(), self._max_extracted_chars)
                    blob.extraction_metadata = metadata
                    blob.extraction_status = "SUCCEEDED"
                except Exception as error:  # Document libraries expose different parse-error types.
                    blob.extraction_status = "FAILED"
                    blob.extraction_error = _bounded(f"{type(error).__name__}: {error}", 1000)
            database.add(blob)
            database.flush()

        linked = EmailAttachmentDb(
            message_id=message.id,
            provider_attachment_id=attachment.attachment_id,
            file_name=file_name,
            blob_hash=blob.content_hash,
        )
        database.add(linked)
        database.flush()
        return linked
